# -*- coding: utf-8 -*-
"""Build the 0529 English progress deck (same template/style as 0522).

Covers: experiment suite (E1-E6) for the conference paper + single-household
DR agent (Phase 2 Step 1, rule-based). Reuses the 0522 deck as theme template.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

PROJECT_ROOT = Path(__file__).resolve().parents[1]      # reproduction/
GRID_DIR = Path(__file__).resolve().parents[2]          # AI Agent smart grid/
NIAR_DIR = Path(__file__).resolve().parents[3]          # 國實院計畫/
ASSETS = PROJECT_ROOT / "slide_assets"
FIGS = PROJECT_ROOT / "figures"
BASE = NIAR_DIR / "0522國實院.pptx"
if not BASE.exists():
    BASE = NIAR_DIR / "0508國實院.pptx"
OUT = NIAR_DIR / "0529國實院.pptx"

NAVY = RGBColor(0x12, 0x41, 0x63); TEAL = RGBColor(0x58, 0xB6, 0xC0)
GREEN = RGBColor(0x75, 0xBD, 0xA7); DARK = RGBColor(0x37, 0x35, 0x45)
GREY = RGBColor(0x7A, 0x8C, 0x8E); LIGHT = RGBColor(0xEA, 0xF1, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); ACCENT = RGBColor(0xE0, 0x7A, 0x5F)
EMU = 914400
EN_FONT = "Corbel"


def _set_font(run, size=None, bold=None, color=None, font=EN_FONT):
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color


def add_slide(prs, i=5): return prs.slides.add_slide(prs.slide_layouts[i])
def set_notes(s, t): s.notes_slide.notes_text_frame.text = t


def set_title(slide, text, color=NAVY, size=27):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Emu(int(0.6*EMU)); ph.top = Emu(int(0.35*EMU))
            ph.width = Emu(int(12.0*EMU)); ph.height = Emu(int(1.0*EMU))
            ph.text = ""; tf = ph.text_frame; tf.word_wrap = True
            try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception: pass
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = text
            _set_font(r, size=size, bold=True, color=color)
            return ph


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(l*EMU)), Emu(int(t*EMU)), Emu(int(w*EMU)), Emu(int(h*EMU)))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf


def add_para(tf, text, size=16, bold=False, color=DARK, bullet=False, space_after=6,
             align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None: p.space_after = Pt(space_after)
    r = p.add_run(); r.text = ("• " + text) if bullet else text
    _set_font(r, size=size, bold=bold, color=color)
    return p


def rect(slide, l, t, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(int(l*EMU)), Emu(int(t*EMU)), Emu(int(w*EMU)), Emu(int(h*EMU)))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def add_image_fit(slide, path, ml, mt, mw, mh):
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = mw/mh
    if ar > bar: w = mw; h = mw/ar
    else: h = mh; w = mh*ar
    slide.shapes.add_picture(str(path), Emu(int((ml+(mw-w)/2)*EMU)), Emu(int((mt+(mh-h)/2)*EMU)),
                             Emu(int(w*EMU)), Emu(int(h*EMU)))


def _table(slide, rows, l, t, w, h, widths, fs=12, header_fs=13):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Emu(int(l*EMU)), Emu(int(t*EMU)),
                                 Emu(int(w*EMU)), Emu(int(h*EMU))).table
    for ci, wd in enumerate(widths):
        tbl.columns[ci].width = Emu(int(wd*EMU))
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci); c.text = ""; c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            if ri == 0:
                _set_font(r, size=header_fs, bold=True, color=WHITE)
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
            else:
                _set_font(r, size=fs, bold=(ci == 0), color=DARK)
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return tbl


def build():
    prs = Presentation(str(BASE))
    lst = prs.slides._sldIdLst
    rel = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for s in list(lst):
        rId = s.get(rel)
        if rId: prs.part.drop_rel(rId)
        lst.remove(s)

    # ---- 1 Title ----
    s = add_slide(prs, 0)
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = ""; p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "NIAR Smart Grid"; _set_font(r, size=40, bold=True, color=NAVY)
            p2 = ph.text_frame.add_paragraph(); r2 = p2.add_run()
            r2.text = "Multiple AI Agents for Residential Energy Management"
            _set_font(r2, size=20, color=TEAL)
        elif idx == 1:
            ph.text = ""; p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "Progress Report ③: Experiment Suite & Single-Household DR Agent"
            _set_font(r, size=17, color=DARK)
            p2 = ph.text_frame.add_paragraph(); r2 = p2.add_run()
            r2.text = "陳柏宇、管少棋   |   2026/05/29"; _set_font(r2, size=14, color=GREY)
    set_notes(s, "Third progress report. Two parts: (A) we completed the full experiment "
                 "suite for the conference paper; (B) we implemented and validated a "
                 "single-household demand-response agent.")

    # ---- 2 Outline ----
    s = add_slide(prs, 5); set_title(s, "This Week")
    items = [
        ("1", "Recap", "Phase 1 done; last week's finding (the 0.94 is an evaluation artifact)"),
        ("2", "Part A — Conference experiment suite (E1–E6)", "granularity, leakage, noise ceiling, vs original method"),
        ("3", "Part B — Single-household DR agent", "design + rule-based controller + results"),
        ("4", "Progress & next steps", "roadmap"),
    ]
    y = 1.7
    for num, t, d in items:
        c = rect(s, 0.8, y, 0.55, 0.55, NAVY); tf = c.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; _set_font(r, size=17, bold=True, color=WHITE)
        tb, tf2 = textbox(s, 1.55, y-0.08, 10.6, 0.85)
        add_para(tf2, t, size=17, bold=True, color=DARK, first=True, space_after=0)
        add_para(tf2, d, size=12, color=GREY, space_after=0)
        y += 1.05
    set_notes(s, "Roadmap of the talk: recap, then the two work streams this week (experiments, agent), then next steps.")

    # ---- 3 Recap ----
    s = add_slide(prs, 5); set_title(s, "Recap: Where We Were")
    tb, tf = textbox(s, 0.7, 1.7, 11.6, 4.6)
    add_para(tf, "Phase 1 (reproduction + forecasting) — done", size=17, bold=True, color=NAVY, first=True, space_after=8)
    for t in ["Reproduced Durrani et al. (2025) on the UCI Appliances dataset",
              "Last week's key finding: the reported R²=0.94 is an EVALUATION ARTIFACT "
              "(coarse-grained evaluation), not a stronger model"]:
        add_para(tf, t, size=14, bullet=True, color=DARK, space_after=7)
    add_para(tf, "This week — two work streams", size=17, bold=True, color=NAVY, space_after=8)
    for t in ["Part A: completed the full experiment suite (E1–E6) that backs the conference paper",
              "Part B: implemented & validated a single-household demand-response agent "
              "(the first step of the multi-agent framework)"]:
        add_para(tf, t, size=14, bullet=True, color=ACCENT, space_after=7)
    set_notes(s, "Remind the professor: Phase 1 is done and last week we showed the 0.94 is "
                 "an evaluation artifact. This week we (A) made that rigorous with a full "
                 "experiment suite, and (B) built the first agent.")

    # ---- 4 Part A overview ----
    s = add_slide(prs, 5); set_title(s, "Part A — Conference Experiment Suite (E1–E6)")
    rows = [
        ["Exp", "Question", "Result"],
        ["E1", "Can we reproduce Durrani's tables?", "Yes (11 models × 7 DR)"],
        ["E2 ★", "Where does 0.94 come from?", "evaluation granularity"],
        ["E3 ★", "Does the data split leak?", "random split inflates R²"],
        ["E4", "vs the original method (Candanedo)?", "reproduce 0.57; ours 0.64"],
        ["E5", "What do lag features do?", "drive gain AND leakage"],
        ["E6 ★", "Where does our model gain come from?", "tuning > CNN/ensemble"],
    ]
    _table(s, rows, 0.8, 1.75, 11.6, 3.6, [1.4, 6.4, 3.8], fs=13)
    tb, tf = textbox(s, 0.8, 5.6, 11.6, 1.0)
    add_para(tf, "Headline: at honest 10-min evaluation every competitive model plateaus at "
                 "R²≈0.64–0.66 (a noise ceiling); higher numbers come only from coarser "
                 "evaluation or split leakage. (Per-experiment figures are in the paper.)",
             size=13, bold=True, color=NAVY, first=True)
    set_notes(s, "Overview table of the six experiments and what each answers. The three "
                 "starred ones are the core contribution. Bottom line is the noise-ceiling message.")

    # ---- 5 Fig A granularity ----
    s = add_slide(prs, 5); set_title(s, "E2 — R² vs Evaluation Granularity (the key finding)")
    add_image_fit(s, FIGS/"figA_granularity.png", 1.2, 1.5, 10.6, 4.3)
    tb, tf = textbox(s, 0.7, 5.9, 11.8, 0.8)
    add_para(tf, "Fig. A.  Same models — coarsening 10-min → 1/2-hour lifts R² from ~0.65 to "
                 "~0.91, approaching the claimed 0.94. The number tracks the evaluation window, not the model.",
             size=12.5, bold=True, color=NAVY, first=True)
    set_notes(s, "The core figure. All models climb from ~0.65 (10-min) to ~0.91 (2-hour) toward "
                 "the 0.94 dashed line. Same predictions, only the scoring window changes. This is "
                 "the mechanism behind the inflated claim.")

    # (E3 leakage / E6 decomposition figures live in the paper — summarized in
    #  the Part-A overview table above; kept out of the weekly deck to stay lean.)

    # ---- 6 Part B: architecture position ----
    s = add_slide(prs, 5); set_title(s, "Part B — Single-Household Agent: Where It Sits")
    add_image_fit(s, ASSETS/"s_full_arch.png", 0.6, 1.45, 11.8, 4.4)
    tb, tf = textbox(s, 0.7, 5.95, 11.8, 0.8)
    add_para(tf, "Data + LSTM Forecast + DR-simulation environment are done; State / Decision / "
                 "Objective now implemented (rule-based). This is one complete household agent — the "
                 "building block of the multi-agent layer.", size=12, bold=True, color=NAVY, first=True)
    set_notes(s, "Place the agent in the full framework. Green = done, teal = rule-based done. We "
                 "completed one full Local Household Intelligence pipeline. Multi-agent / global / "
                 "federated remain future work (and we already have a multi-household dataset for that).")

    # ---- 9 Agent method ----
    s = add_slide(prs, 5); set_title(s, "Agent Method — Deferrable-Load Decision Loop")
    add_image_fit(s, ASSETS/"s_agent_loop.png", 0.6, 1.4, 11.8, 3.4)
    tb, tf = textbox(s, 0.7, 5.0, 11.8, 1.7)
    add_para(tf, "Since the dataset is a static log (no price, no actions), we wrap it in a "
                 "simulator: 30% of each step's load is flexible and can be deferred into a buffer, "
                 "auto-released off-peak. The agent defers when price is high and a peak is forecast.",
             size=13, color=DARK, first=True, space_after=6)
    add_para(tf, "Loop:  predict (LSTM)  →  build state  →  decide (rule)  →  act (defer/release) "
                 " →  environment reshapes load  →  reward.   Concrete setup on the next slide.",
             size=12.5, bold=True, color=NAVY)
    set_notes(s, "Explain the loop left to right: forecast → state → rule-based decision → action "
                 "(defer/release) → environment reshapes load → reward. Stress the honest caveats: "
                 "synthetic price, assumed 30% flexible fraction, auto-release rule.")

    # ---- 10 Agent setup (concrete configuration) ----
    s = add_slide(prs, 5); set_title(s, "Agent Setup — State / Action / Reward / Environment")

    def _titled_box(x, y, w, h, title, lines, tcol=NAVY):
        b = rect(s, x, y, w, h, LIGHT)
        tf = b.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(12); tf.margin_right = Pt(10); tf.margin_top = Pt(8)
        add_para(tf, title, size=14, bold=True, color=tcol, first=True, space_after=5)
        for ln in lines:
            add_para(tf, ln, size=11.5, color=DARK, bullet=True, space_after=3)

    _titled_box(0.7, 1.6, 5.7, 2.35, "Environment (deferrable-load simulator)", [
        "Flexible fraction φ = 30% of each step's load",
        "Buffer cap 600 Wh; release ≤ 150 Wh/step (off-peak)",
        "Synthetic ToU price: peak(17–22)=0.30, off(0–7)=0.08, mid=0.15",
    ])
    _titled_box(6.6, 1.6, 5.9, 2.35, "State  (6-dim)", [
        "current demand · LSTM next-step forecast",
        "current price · hour sin/cos",
        "buffer level (deferred energy queued)",
    ])
    _titled_box(0.7, 4.15, 5.7, 1.9, "Action  (discrete, 3)", [
        "0 = serve normally",
        "1 = defer half of flexible load",
        "2 = defer all flexible load",
    ])
    _titled_box(6.6, 4.15, 5.9, 1.9, "Reward", [
        "R = −(w₁·cost + w₂·peak + w₃·comfort + w₄·switching)",
        "weights 1.0 / 1.0 / 0.3 / 0.05",
        "cost=price·served · comfort=undelivered buffer",
    ])
    tb, tf = textbox(s, 0.7, 6.2, 11.8, 0.7)
    add_para(tf, "Dynamics: served = demand − deferred + released.   "
                 "Rule-based policy: high price + forecast peak → defer; off-peak → auto-release.",
             size=12, bold=True, color=ACCENT, first=True)
    set_notes(s, "This is the concrete agent configuration. Environment: 30% of load is flexible, "
                 "stored in a 600 Wh buffer, released off-peak; synthetic Time-of-Use price. State is "
                 "a 6-dim vector (demand, forecast, price, time, buffer). Action is discrete 3 "
                 "(defer none/half/all). Reward penalizes cost, peak, undelivered energy (comfort), "
                 "and switching. Served load = demand − deferred + released. Be explicit that φ and "
                 "the price are assumptions to be sensitivity-tested.")

    # ---- 11 Agent results ----
    s = add_slide(prs, 5); set_title(s, "Agent Results — vs No-DR Baseline")
    add_image_fit(s, FIGS/"figF_agent_loadcurve.png", 0.5, 1.5, 7.4, 3.6)
    rows = [
        ["Metric", "Baseline → Agent"],
        ["Electricity cost", "−2.6%"],
        ["Peak-window avg (17–22)", "−5.6%"],
        ["95th-pctile load (peak shaving)", "−10.3%"],
        ["Energy served vs demand", "conserved ✓"],
    ]
    _table(s, rows, 8.1, 1.7, 4.3, 2.6, [2.7, 1.6], fs=11, header_fs=12)
    tb, tf = textbox(s, 8.1, 4.5, 4.3, 2.0)
    add_para(tf, "The agent shaves evening peaks and shifts to off-peak — cost down, peak down, "
                 "energy conserved. Loop validated end-to-end.", size=12, bold=True, color=ACCENT, first=True)
    tb2, tf2 = textbox(s, 0.5, 5.25, 7.4, 0.7)
    add_para(tf2, "Fig. F.  Blue (agent) clipped below the threshold at peaks vs grey (baseline).",
             size=11, color=GREY, first=True)
    set_notes(s, "Results: cost −2.6%, peak-window −5.6%, 95th-percentile load −10.3%, energy "
                 "conserved (no cheating by under-serving). The figure shows blue agent load clipped "
                 "below the grey baseline at peaks. The predict→decide→act→reward loop is validated.")

    # ---- 11 Future direction: LLM advisory layer ----
    s = add_slide(prs, 5); set_title(s, "Future Direction — LLM Advisory Layer")
    flow = [
        ("Forecast", "LSTM next-hour\nload / peak", GREEN, WHITE),
        ("Decision Core", "rule → MPC\n(optimal DR)", TEAL, WHITE),
        ("LLM Advisory", "personalized,\nexplainable advice", NAVY, WHITE),
        ("User", "report +\nrecommendation", ACCENT, WHITE),
    ]
    bw = 2.5; bh = 1.5; gap = 0.55; x0 = (13.3 - (4*bw + 3*gap))/2; y = 1.9
    from pptx.enum.shapes import MSO_SHAPE
    for i,(t, sub, fc, tc) in enumerate(flow):
        x = x0 + i*(bw+gap)
        b = rect(s, x, y, bw, bh, fc)
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t; _set_font(r, size=13, bold=True, color=tc)
        add_para(tf, sub, size=10, color=tc, align=PP_ALIGN.CENTER, space_after=0)
        if i < 3:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int((x+bw+0.05)*EMU)),
                                    Emu(int((y+bh/2-0.12)*EMU)), Emu(int((gap-0.1)*EMU)),
                                    Emu(int(0.24*EMU)))
            ar.fill.solid(); ar.fill.fore_color.rgb = GREY; ar.line.fill.background()
    tb, tf = textbox(s, 0.7, 3.9, 11.9, 2.7)
    add_para(tf, "Why advisory, not direct device control", size=15, bold=True, color=NAVY,
             first=True, space_after=7)
    for t in ["Direct device control is hard in practice (hardware integration, protocols, "
              "safety, user trust) — advisory is deployable today",
              "LLM turns the optimizer's output into personalized, natural-language "
              "recommendations + explanations (directly targets the XAI / trust gap)",
              "Guardrail: the LLM only phrases/explains; all numbers come from the decision core",
              "Hybrid: operate in advisory mode by default; optional direct control where hardware allows",
              "Evaluation: 'potential savings if followed' (our simulator already computes this) "
              "+ explanation faithfulness"]:
        add_para(tf, t, size=12.5, bullet=True, color=DARK, space_after=5)
    set_notes(s, "Future positioning. Direct control of real devices is hard, so we position the "
                 "system as ADVISORY: forecast → decision core computes optimal DR → an LLM turns it "
                 "into personalized, explainable recommendations for the user. The LLM is the "
                 "explanation/interface layer, not the decision logic (guardrail: it only phrases, "
                 "doesn't compute). This targets the XAI/trust gap the literature highlights, is "
                 "deployable, and doesn't waste our forecast+decision work. Hybrid: optional direct "
                 "control where hardware allows.")

    # ---- 12 Roadmap ----
    s = add_slide(prs, 5); set_title(s, "Progress & Next Steps")
    rows = [
        ["Stage", "Status"],
        ["Phase 1 — reproduction + experiment suite (E1–E6)", "Done"],
        ["Phase 2 Step 1 — single-household agent (rule-based)", "Done (this week)"],
        ["Phase 2 Step 1b — MPC baseline / optional RL", "Next"],
        ["Phase 2 Step 2 — LLM advisory layer (report + explanation)", "Planned"],
        ["Phase 3 — multi-agent coordination (multi-household data)", "Planned"],
        ["Phase 3 — federated learning + global coordination", "Future"],
    ]
    _table(s, rows, 0.9, 1.8, 11.4, 3.5, [8.6, 2.8], fs=12)
    tb, tf = textbox(s, 0.9, 5.6, 11.4, 1.1)
    add_para(tf, "Decision: pause at the rule-based agent (loop validated). Conference paper uses "
                 "Phase 1 (E1–E6); the agent extends to multi-household + an LLM advisory layer next.",
             size=13, bold=True, color=ACCENT, first=True)
    set_notes(s, "Roadmap. Phase 1 and Phase 2 Step 1 done. Next: MPC baseline, then the LLM "
                 "advisory layer, then multi-agent coordination on the multi-household dataset.")

    # ---- 12 References ----
    s = add_slide(prs, 5); set_title(s, "References")
    refs = [
        ("[1]", "Durrani A. M., et al. (2025). AI-driven optimization of energy consumption and "
                "demand response in smart homes. Energy Exploration & Exploitation, 44(3).",
                "DOI: 10.1177/01445987251403607  — paper reproduced"),
        ("[2]", "Candanedo L. M., et al. (2017). Data driven prediction models of energy use of "
                "appliances in a low-energy house. Energy and Buildings, 140.",
                "DOI: 10.1016/j.enbuild.2017.01.083  — dataset origin; GBM test R²≈0.57"),
        ("[3]", "Kulkarni P. (2025). Appliance Energy Prediction using Machine Learning Techniques. "
                "Master's thesis, MTSU.", "LSTM 0.60 / GRU 0.62"),
        ("[4]", "Dataset: UCI Appliances Energy Prediction.",
                "archive.ics.uci.edu/dataset/374"),
    ]
    y = 1.75
    for tag, cite, meta in refs:
        tb, tf = textbox(s, 0.8, y, 11.7, 1.0)
        p = tf.paragraphs[0]; p.space_after = Pt(2)
        r = p.add_run(); r.text = tag+"  "; _set_font(r, size=12.5, bold=True, color=NAVY)
        r2 = p.add_run(); r2.text = cite; _set_font(r2, size=12.5, color=DARK)
        add_para(tf, meta, size=10.5, color=GREY, space_after=0)
        y += 1.12
    set_notes(s, "References. [1] reproduced paper, [2] dataset origin (0.57 benchmark), [3] "
                 "deep-learning comparison, [4] dataset link.")

    # ---- 13 Thanks ----
    s = add_slide(prs, 0)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""; p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "Thanks for listening!"; _set_font(r, size=40, bold=True, color=NAVY)
        elif ph.placeholder_format.idx == 1:
            ph.text = ""; p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "Q & A"; _set_font(r, size=20, color=TEAL)
    set_notes(s, "Likely Q&A: (1) why not RL now? — rule-based validates the loop; MPC is the "
                 "principled next baseline; RL needs a credible (non-synthetic) reward. (2) single "
                 "household value? — it's the building block; grid value comes from multi-agent "
                 "coordination, next phase with the multi-household dataset.")

    prs.save(str(OUT))
    print("saved ->", OUT, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
