# -*- coding: utf-8 -*-
"""Build the 0522 English progress deck from the 0515 template.

Keeps the 0515 theme (colors / fonts / layouts), removes its slides, adds new
English slides with paper-style Fig./Table. numbering, embeds speaker notes.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# reports/build_0522.py -> parents[1] = project root (reproduction/),
# parents[2] = "AI Agent smart grid" (where the weekly .pptx decks live)
PROJECT_ROOT = Path(__file__).resolve().parents[1]
GRID_DIR = Path(__file__).resolve().parents[2]
BASE = GRID_DIR / "0515國實院.pptx"     # NOTE: template was removed; restore if rebuilding
ASSETS = PROJECT_ROOT / "slide_assets"
OUT = GRID_DIR / "0522國實院.pptx"

NAVY = RGBColor(0x12, 0x41, 0x63)
TEAL = RGBColor(0x58, 0xB6, 0xC0)
GREEN = RGBColor(0x75, 0xBD, 0xA7)
DARK = RGBColor(0x37, 0x35, 0x45)
GREY = RGBColor(0x7A, 0x8C, 0x8E)
LIGHT = RGBColor(0xEA, 0xF1, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE0, 0x7A, 0x5F)

EMU = 914400
EN_FONT = "Corbel"


def _set_font(run, size=None, bold=None, color=None, font=EN_FONT):
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color


def add_slide(prs, layout_idx=5):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def set_title(slide, text, color=NAVY, size=28):
    title = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title = ph; break
    if title is None:
        return None
    title.left = Emu(int(0.6*EMU)); title.top = Emu(int(0.35*EMU))
    title.width = Emu(int(12.0*EMU)); title.height = Emu(int(1.0*EMU))
    title.text = ""
    tf = title.text_frame; tf.word_wrap = True
    try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception: pass
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    _set_font(r, size=size, bold=True, color=color)
    return title


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(l*EMU)), Emu(int(t*EMU)),
                                  Emu(int(w*EMU)), Emu(int(h*EMU)))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf


def add_para(tf, text, size=16, bold=False, color=DARK, bullet=False,
             space_after=6, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None: p.space_after = Pt(space_after)
    r = p.add_run(); r.text = ("• " + text) if bullet else text
    _set_font(r, size=size, bold=bold, color=color)
    return p


def rect(slide, l, t, w, h, fill, round=True):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Emu(int(l*EMU)), Emu(int(t*EMU)), Emu(int(w*EMU)), Emu(int(h*EMU)))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def add_image_fit(slide, path, max_l, max_t, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih; box_ar = max_w/max_h
    if ar > box_ar: w = max_w; h = max_w/ar
    else: h = max_h; w = max_h*ar
    l = max_l + (max_w-w)/2; t = max_t + (max_h-h)/2
    slide.shapes.add_picture(str(path), Emu(int(l*EMU)), Emu(int(t*EMU)),
                             Emu(int(w*EMU)), Emu(int(h*EMU)))


def build():
    prs = Presentation(str(BASE))
    sldIdLst = prs.slides._sldIdLst
    rel_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sld in list(sldIdLst):
        rId = sld.get(rel_attr)
        if rId: prs.part.drop_rel(rId)
        sldIdLst.remove(sld)

    # ---------- Slide 1: Title ----------
    s = add_slide(prs, 0)
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "NIAR Smart Grid"
            _set_font(r, size=40, bold=True, color=NAVY)
            p2 = ph.text_frame.add_paragraph()
            r2 = p2.add_run(); r2.text = "Multiple AI Agents for Residential Energy Management"
            _set_font(r2, size=20, color=TEAL)
        elif idx == 1:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "Reproduction Progress Report ②: LSTM Improvement & Evaluation Analysis"
            _set_font(r, size=17, color=DARK)
            p2 = ph.text_frame.add_paragraph()
            r2 = p2.add_run(); r2.text = "陳柏宇、管少棋   |   2026/05/22"
            _set_font(r2, size=14, color=GREY)
    set_notes(s, "Opening. This is our second reproduction report. Last week we showed "
                 "the honest replication and the gap to the paper; this week we (1) explain "
                 "WHY the paper reaches 0.94, (2) improve our own LSTM, and (3) place the work "
                 "in the full system architecture.")

    # ---------- Slide 2: Outline ----------
    s = add_slide(prs, 5)
    set_title(s, "This Week")
    items = [
        ("1", "Phase 1 reproduction recap", "7-step pipeline · 8 models × 7 DR strategies"),
        ("2", "Replication results: honest vs paper", "the real gap under 10-min evaluation"),
        ("3", "Key finding: evaluation granularity", "0.94 comes from coarser-time evaluation"),
        ("4", "Literature comparison", "the realistic R² range for this dataset"),
        ("5", "LSTM improvement: last week vs this week", "CNN-LSTM + cyclical features + ensemble"),
        ("6", "Improved full results", "raw / MA6 / MA12 vs paper"),
        ("7", "System architecture & current stage", "where we are in the framework"),
        ("8", "Progress table & next steps", ""),
    ]
    y = 1.45
    for num, title, desc in items:
        c = rect(s, 0.7, y, 0.5, 0.5, NAVY)
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; _set_font(r, size=16, bold=True, color=WHITE)
        tb, tf2 = textbox(s, 1.4, y-0.06, 10.6, 0.66)
        add_para(tf2, title, size=16, bold=True, color=DARK, first=True, space_after=0)
        if desc:
            add_para(tf2, desc, size=11, color=GREY, space_after=0)
        y += 0.66
    set_notes(s, "Walk through the 8 sections in one breath. Emphasize that the storyline is: "
                 "we reproduced it, we explain the suspicious 0.94, we genuinely improved our model, "
                 "and we show the roadmap.")

    # ---------- Slide 3: Phase 1 recap ----------
    s = add_slide(prs, 5)
    set_title(s, "Phase 1 Reproduction Recap: Methods & Outputs")
    tb, tf = textbox(s, 0.6, 1.55, 6.1, 5.0)
    add_para(tf, "Seven-step pipeline fully reproduced", size=17, bold=True, color=NAVY,
             first=True, space_after=9)
    for t in ["Data preprocessing (impute / dedup / 3σ outliers / Min-Max)",
              "Feature engineering (time, apparent temp, discomfort index, lag features)",
              "EDA (time series, correlation heatmap, histogram, pair plot)",
              "8 models (LR / RF / SVR / kNN / LSTM + 3 baselines)",
              "7 DR strategies",
              "Evaluation: MAE / RMSE / R²  (8 × 7 = 56 cells)",
              "Visualization: heatmaps, bar plots, prediction lines"]:
        add_para(tf, t, size=12.5, bullet=True, color=DARK, space_after=6)
    box = rect(s, 7.0, 1.6, 5.4, 3.3, LIGHT)
    tf = box.text_frame; tf.word_wrap = True; tf.margin_left=Pt(14); tf.margin_top=Pt(12)
    add_para(tf, "Aligned with the paper, item by item", size=15, bold=True, color=NAVY,
             first=True, space_after=8)
    for t in ["Table 2 descriptive stats: exact match\n(mean 97.69 / std 102.52 / max 1080)",
              "Model ranking trend matches:\nLSTM ≈ RF > kNN > SVR > LR",
              "Naive baselines correctly collapse to\nR² = −1 on Load Leveling / Shifting"]:
        add_para(tf, t, size=12, bullet=True, color=DARK, space_after=10)
    tb, tf = textbox(s, 7.0, 5.05, 5.4, 1.1)
    add_para(tf, "Conclusion: methodology 100% reproduced;\nnumeric gap analyzed next.",
             size=13, bold=True, color=ACCENT, first=True)
    set_notes(s, "Say: 'Phase 1 is fully reproduced.' Point to the LEFT list as the 7 pipeline "
                 "steps from the paper. Then point to the RIGHT box: the three things that prove "
                 "our reproduction is faithful — the descriptive statistics match the paper exactly, "
                 "the model ranking is the same, and the naive baselines fail in exactly the way the "
                 "paper reports. Finish with the orange line: method is reproduced, the number gap is "
                 "what the rest of the talk explains.")

    # ---------- Slide 4: Table 1 ----------
    s = add_slide(prs, 5)
    set_title(s, "Replication Results (Honest): Ours vs Durrani et al. (2025)")
    tb, tf = textbox(s, 0.9, 1.55, 11, 0.5)
    add_para(tf, "Table 1.  LSTM headline (Price-Based DR), raw 10-min evaluation.",
             size=14, bold=True, color=DARK, first=True)
    rows = [
        ("Metric", "Ours (honest)", "Paper (claimed)", "Gap"),
        ("R²", "0.59 – 0.64", "0.94", "−0.30"),
        ("MAE (Wh)", "≈ 23", "18.95", "+4"),
        ("RMSE (Wh)", "≈ 46", "24.83", "+21"),
    ]
    tbl = s.shapes.add_table(len(rows), 4, Emu(int(0.9*EMU)), Emu(int(2.2*EMU)),
                             Emu(int(10.4*EMU)), Emu(int(2.1*EMU))).table
    for ci, wd in enumerate([2.6, 3.0, 2.6, 2.2]):
        tbl.columns[ci].width = Emu(int(wd*EMU))
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            if ri == 0:
                _set_font(r, size=14, bold=True, color=WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                _set_font(r, size=13.5, bold=(ci==0), color=DARK)
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri%2 else LIGHT
    tb, tf = textbox(s, 0.9, 4.7, 10.6, 1.7)
    add_para(tf, "MAE is actually close — the gap is in R² and RMSE.", size=14, bold=True,
             color=NAVY, first=True, space_after=6)
    add_para(tf, "Clue: our RMSE/MAE ≈ 2.0 (heavy-tail errors) vs the paper's ≈ 1.3 (near-Gaussian)"
                 " — the paper's few high-residual peak points have been 'smoothed away'.",
             size=12.5, color=DARK)
    set_notes(s, "This is Table 1. Read the three rows. The key sentence: 'Our average error (MAE) is "
                 "actually close to the paper — only about 4 Wh higher. The big gap is in R² and RMSE.' "
                 "Then explain the clue: RMSE divided by MAE. For us it's about 2.0, which means we have "
                 "a few very large errors (the unpredictable spikes). For the paper it's about 1.3, which "
                 "is what you get when those spikes are gone. That points to smoothing — which is the next slide.")

    # ---------- Slide 5: key finding ----------
    s = add_slide(prs, 5)
    set_title(s, "Key Finding: the R² Gap Comes from Evaluation Granularity")
    tb, tf = textbox(s, 0.7, 1.55, 11.5, 1.0)
    add_para(tf, "Same model, same predictions — switching from 'point-by-point at 10 min' to "
                 "'hourly average' lifts R² from 0.6 to 0.88+.", size=15, bold=True,
             color=DARK, first=True)
    b1 = rect(s, 0.7, 2.7, 5.4, 3.2, LIGHT)
    tf = b1.text_frame; tf.word_wrap=True; tf.margin_left=Pt(14); tf.margin_top=Pt(12)
    add_para(tf, "Raw 10-min evaluation", size=16, bold=True, color=GREY, first=True, space_after=8)
    for t in ["Must predict 'is the hair-dryer on at exactly 18:23?'",
              "Instantaneous spikes are random & unpredictable",
              "R² ≈ 0.59 – 0.64"]:
        add_para(tf, t, size=13, bullet=True, color=DARK, space_after=8)
    b2 = rect(s, 6.5, 2.7, 5.4, 3.2, GREEN)
    tf = b2.text_frame; tf.word_wrap=True; tf.margin_left=Pt(14); tf.margin_top=Pt(12)
    add_para(tf, "Hourly (MA6 / MA12) evaluation", size=16, bold=True, color=WHITE, first=True, space_after=8)
    for t in ["Only predict 'total energy this hour'",
              "Random spikes cancel out — only the pattern is judged",
              "R² ≈ 0.85 – 0.91"]:
        add_para(tf, t, size=13, bullet=True, color=WHITE, space_after=8)
    tb, tf = textbox(s, 0.7, 6.1, 11.5, 0.7)
    add_para(tf, "Neutral conclusion: the paper most likely computes its metrics at a coarser time "
                 "granularity (not stated in the text) — not a fundamentally stronger model.",
             size=13, bold=True, color=NAVY, first=True)
    set_notes(s, "This is the heart of the talk. Use the hair-dryer analogy. LEFT box: at 10-minute "
                 "resolution the model must guess whether someone happens to switch on a high-power "
                 "appliance at a precise instant — that is essentially random, so R² stays ~0.6. "
                 "RIGHT box: at hourly resolution the random spikes average out, the model only has to "
                 "get the overall hourly level right, which is easy, so R² jumps to ~0.9. Stress the "
                 "neutral wording at the bottom: we are NOT accusing them of cheating — we are saying "
                 "the evaluation granularity differs, and it is not documented.")

    # ---------- Slide 6: Fig 1 literature ----------
    s = add_slide(prs, 5)
    set_title(s, "Literature Comparison: Realistic R² Range for this Dataset")
    add_image_fit(s, ASSETS/"s_literature.png", 0.7, 1.45, 11.5, 4.1)
    tb, tf = textbox(s, 0.7, 5.7, 11.6, 1.0)
    add_para(tf, "Fig. 1.  Candanedo (2017) GBM = 0.57; lit. LSTM = 0.60; GRU = 0.62. "
                 "Our honest raw R² = 0.64 sits at the top of the literature; 0.94 is a clear outlier.",
             size=12.5, bold=True, color=NAVY, first=True)
    set_notes(s, "Fig. 1. This slide defends our honest number with the literature. The original "
                 "dataset paper (Candanedo 2017) got 0.57 with gradient boosting. Follow-up deep "
                 "learning studies report 0.60–0.62. Our improved LSTM at 0.64 (grey-green bar) is "
                 "already at the TOP of everything published at 10-min resolution. The point: 0.94 is "
                 "an outlier that nobody else reaches honestly — which supports the granularity "
                 "explanation. The teal bar (0.91) is OUR hourly-evaluation number, right next to the paper.")

    # ---------- Slide 7: Fig 2 week-over-week ----------
    s = add_slide(prs, 5)
    set_title(s, "LSTM Improvement: Last Week vs This Week")
    add_image_fit(s, ASSETS/"s_week_over_week.png", 0.5, 1.5, 5.9, 4.6)
    tb, tf = textbox(s, 6.6, 1.55, 5.5, 5.0)
    add_para(tf, "What we changed this week", size=16, bold=True, color=NAVY, first=True, space_after=9)
    for t in ["CNN-LSTM hybrid (Conv1D for local patterns)",
              "Cyclical time features (sin/cos of hour & day)",
              "4-model ensemble (multi-seed averaging)",
              "Back to MSE loss, removed log transform",
              "Validation-based early stopping"]:
        add_para(tf, t, size=13, bullet=True, color=DARK, space_after=8)
    add_para(tf, "Key lesson", size=14.5, bold=True, color=ACCENT, space_after=5)
    add_para(tf, "log transform + Huber loss suppress peak prediction —\nbetter MAE but worse R²/RMSE, so we removed them.",
             size=12, color=DARK)
    set_notes(s, "Fig. 2. This is the genuine model-improvement slide. Compared with last week's "
                 "0.579, this week's LSTM reaches 0.644 on the SAME honest 10-min evaluation — a real "
                 "+0.065 gain, purely from architecture, not from changing the metric. List the five "
                 "changes on the right. End with the lesson: we first tried log-transform and Huber "
                 "loss; they improved MAE but HURT R², because they stop the model from chasing peaks. "
                 "We reverted them. This shows we understand the metric, not just tuning blindly.")

    # ---------- Slide 8: Fig 3 v2 vs paper ----------
    s = add_slide(prs, 5)
    set_title(s, "Improved Full Results: raw / MA6 / MA12 vs Paper")
    add_image_fit(s, ASSETS/"s_v2_vs_paper.png", 0.6, 1.5, 11.6, 4.2)
    tb, tf = textbox(s, 0.7, 5.75, 11.6, 1.0)
    add_para(tf, "Fig. 3.  MA12 (2 h) mean R² = 0.909 vs paper 0.921 — essentially matched; "
                 "peak_clipping and load_leveling now exceed the paper.",
             size=12.5, bold=True, color=NAVY, first=True)
    set_notes(s, "Fig. 3 — the full results. "
                 "Four bars per DR strategy: grey = our raw 10-min, green = 1-hour eval, teal = 2-hour "
                 "eval, navy = the paper. Walk left to right: raw bars are 0.6–0.79; as we coarsen the "
                 "evaluation the green and teal bars climb to meet the navy paper bar. At 2-hour "
                 "granularity our average (0.909) basically equals the paper (0.921), and for peak "
                 "clipping and load leveling we even exceed it. Conclusion: once the evaluation matches, "
                 "our model is on par with the paper.")

    # ---------- Slide 9: Fig 4 full architecture ----------
    s = add_slide(prs, 5)
    set_title(s, "System Architecture & Current Stage (0508 p.6)")
    add_image_fit(s, ASSETS/"s_full_arch.png", 0.5, 1.45, 11.9, 4.4)
    tb, tf = textbox(s, 0.6, 5.95, 12.0, 0.9)
    add_para(tf, "Fig. 4.  Completed: Data + LSTM Forecast (Phase 1). DR Simulation is rule-based "
                 "(not yet agent decision). Global coordination, State / Decision / Objective, and "
                 "device execution remain future work.", size=12, bold=True, color=NAVY, first=True)
    set_notes(s, "Fig. 4 — the big picture. This is the full architecture from the 0508 proposal. "
                 "Three layers: a global coordination layer on top, the household-agent layer in the "
                 "middle, physical devices and a feedback loop at the bottom. Inside each household is "
                 "the Local Household Intelligence pipeline. The green boxes (Data, LSTM Forecast) are "
                 "DONE — that is the orange 'You are here, Phase 1 complete' arrow. The teal box (DR "
                 "Simulation) is partial and rule-based. Everything grey — State, Decision Core, "
                 "Objective, the global layer, device execution — is future work. So we have finished "
                 "the bottom-left of the perception/forecasting part and are about to enter decision-making.")

    # ---------- Slide 10: Fig 5 local intelligence ----------
    s = add_slide(prs, 5)
    set_title(s, "Zoom-in: Local Household Intelligence Pipeline")
    add_image_fit(s, ASSETS/"s_arch_stage.png", 0.5, 1.7, 11.9, 3.4)
    tb, tf = textbox(s, 0.6, 5.3, 12.0, 1.3)
    add_para(tf, "Fig. 5.  We have Data + LSTM Forecast; DR Simulation is rule-based, not yet an "
                 "agent decision. Immediate milestone: complete ONE full single-household agent "
                 "(State → Decision Core → Objective) before scaling to multi-agent / federated learning.",
             size=13, bold=True, color=NAVY, first=True)
    set_notes(s, "Fig. 5 — the zoom-in of the per-household pipeline. Restate clearly which stages are "
                 "done (Data, LSTM Forecast = green), partial (DR Simulation = teal), and not started "
                 "(State, Decision Core, Objective = grey). The transition we are about to make is from "
                 "PREDICTION to DECISION: right now we only forecast load; next we need to build the "
                 "state representation, the decision agent, and the objective function.")

    # ---------- Slide 11: Table 2 progress ----------
    s = add_slide(prs, 5)
    set_title(s, "Progress Table & Next Steps")
    tb, tf = textbox(s, 0.7, 1.5, 11, 0.45)
    add_para(tf, "Table 2.  Roadmap — single-household agent first, then scale outward.",
             size=13.5, bold=True, color=DARK, first=True)
    rows = [
        ("Stage", "Task", "Status"),
        ("This week", "LSTM v2 improvement · evaluation-granularity analysis · literature comparison\n"
                      "(carry-over: integrate v2 into full heatmaps, add chronological-split honest version)", "Done"),
        ("Step 1 —\nSingle-household agent", "Define State / Action / Reward schema · build Decision Core for ONE household "
                      "(rule/workflow → RL) · close the forecast → decision → DR action → reward loop", "Next"),
        ("Step 2 —\nMulti-agent", "Replicate the validated single agent to N households · add coordination "
                      "to avoid synchronized off-peak rebound (new peaks)", "Planned"),
        ("Step 3 —\nFederated + global", "Federated learning across households · global grid-coordination layer · "
                      "Workflow vs Multi-Agent vs Hybrid · grid-scenario evaluation", "Future"),
    ]
    tbl = s.shapes.add_table(len(rows), 3, Emu(int(0.6*EMU)), Emu(int(1.95*EMU)),
                             Emu(int(12.0*EMU)), Emu(int(3.5*EMU))).table
    for ci, wd in enumerate([2.7, 7.1, 2.2]):
        tbl.columns[ci].width = Emu(int(wd*EMU))
    statc = {"Done": GREEN, "Next": TEAL, "Planned": NAVY, "Future": GREY}
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci != 1 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            if ri == 0:
                _set_font(r, size=13, bold=True, color=WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                col = statc.get(val, DARK) if ci == 2 else (NAVY if ci == 0 else DARK)
                _set_font(r, size=10.5, bold=(ci != 1), color=col)
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri%2 else LIGHT
    tb, tf = textbox(s, 0.6, 5.7, 12.0, 1.1)
    add_para(tf, "Strategy: finish and validate ONE complete household agent first — only then "
                 "replicate to multi-agent, and finally add federated learning + global coordination. "
                 "Build the single unit before the system.",
             size=12.5, bold=True, color=ACCENT, first=True)
    set_notes(s, "Table 2 — the roadmap, and the key message is the sequencing. This week (done): model "
                 "improvement + granularity analysis; plus two small carry-over engineering tasks. "
                 "Step 1 is our immediate focus: build ONE complete single-household agent — define the "
                 "State/Action/Reward, build the Decision Core (start rule-based, then reinforcement "
                 "learning), and close the full loop from forecast to decision to action to reward. "
                 "Only after that single agent works do we go to Step 2, multi-agent: replicate it to "
                 "many households and add coordination so they don't all shift load to the same off-peak "
                 "hour and create a NEW peak. Step 3 is federated learning plus the global coordination "
                 "layer, and comparing Workflow / Multi-Agent / Hybrid designs. Bottom line: build and "
                 "validate the single unit before scaling to the system.")

    # ---------- Slide 12: References ----------
    s = add_slide(prs, 5)
    set_title(s, "References")
    refs = [
        ("[1]", "Durrani A. M., et al. (2025). AI-driven optimization of energy consumption and "
                "demand response in smart homes. Energy Exploration & Exploitation, 44(3), 1382–1419.",
                "DOI: 10.1177/01445987251403607   — the paper reproduced in this work"),
        ("[2]", "Candanedo L. M., Feldheim V., & Deramaix D. (2017). Data driven prediction models "
                "of energy use of appliances in a low-energy house. Energy and Buildings, 140, 81–97.",
                "DOI: 10.1016/j.enbuild.2017.01.083   — dataset origin; GBM testing R² ≈ 0.57"),
        ("[3]", "Kulkarni P. (2025). Appliance Energy Prediction using Machine Learning Techniques. "
                "Master's thesis, Middle Tennessee State University.",
                "jewlscholar.mtsu.edu/items/918b968e-4c23-465b-81fd-376258d21609   — LR 0.19 / LSTM 0.60 / GRU 0.62 (validation)"),
        ("[4]", "Dataset: UCI Appliances Energy Prediction.",
                "archive.ics.uci.edu/dataset/374/appliances+energy+prediction"),
    ]
    y = 1.7
    for tag, cite, meta in refs:
        tb, tf = textbox(s, 0.8, y, 11.6, 1.0)
        p = tf.paragraphs[0]; p.space_after = Pt(2)
        r = p.add_run(); r.text = tag + "  "; _set_font(r, size=13, bold=True, color=NAVY)
        r2 = p.add_run(); r2.text = cite; _set_font(r2, size=13, color=DARK)
        add_para(tf, meta, size=10.5, color=GREY, space_after=0)
        y += 1.18
    set_notes(s, "References. [1] is the paper we reproduce. [2] is the dataset's original paper, "
                 "the source of the 0.57 benchmark. [3] is the deep-learning comparison giving the "
                 "0.60 / 0.62 numbers. [4] is the dataset link. Full source notes with screenshot "
                 "locations are in 0522_文獻來源.md.")

    # ---------- Slide 13: thanks ----------
    s = add_slide(prs, 0)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "Thanks for listening!"
            _set_font(r, size=40, bold=True, color=NAVY)
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "Q & A"
            _set_font(r, size=20, color=TEAL)
    set_notes(s, "Thank the audience. Be ready for: (1) 'Are you sure the paper smoothed?' — answer: "
                 "we cannot confirm without their code, but it is the only assumption that reproduces "
                 "0.94 without leakage. (2) 'What is the chronological-split number?' — about 0.45–0.55, "
                 "honest worst case. (3) 'What is next?' — State/Action/Reward and the decision agent.")

    prs.save(str(OUT))
    print("saved ->", OUT, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
